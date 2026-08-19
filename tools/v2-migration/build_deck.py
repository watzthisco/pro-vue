#!/usr/bin/env python3
"""Build Professional-Vue-v2.0.pptx from the v1.5 deck."""
import os, shutil, sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
import pptxlib as P
import slidecontent as C

ROOT = '/home/user/pro-vue'
SRC = os.path.join(ROOT, 'presentation/Version 1.5/Professional-Vue-v1.5.pptx')
DST = os.path.join(ROOT, 'presentation/Version 2.0/Professional-Vue-v2.0.pptx')

os.makedirs(os.path.dirname(DST), exist_ok=True)
shutil.copyfile(SRC, DST)

prs = Presentation(DST)
original_count = len(prs.slides)

# ---- 1. rewrite text on existing slides (indices are still the originals) ----
rewritten = 0
problems = []
for num, spec in sorted(C.REWRITE.items()):
    slide = prs.slides[num - 1]
    try:
        P.set_slide(slide, title=spec.get('title'), body=spec.get('body'))
        rewritten += 1
    except ValueError:
        # The content placeholder holds a picture or a diagram rather than
        # text -- replace it with a text box in the same place.
        if spec.get('title'):
            P.set_slide(slide, title=spec['title'])
        if spec.get('body') and P.replace_graphic_with_text(slide, spec['body']):
            rewritten += 1
        else:
            problems.append(f'slide {num}: no body placeholder')
print(f'rewrote {rewritten} slides')
for p in problems:
    print('  PROBLEM', p)

# ---- 2. the title slide's subtitle ----
sub = None
for sh in prs.slides[0].shapes:
    if sh.has_text_frame and 'Subtitle' in sh.name:
        sub = sh
if sub is not None:
    P.set_text_frame(sub, [(0, 'copyright 2026, Chris Minnick'),
                           (0, 'version 2.0, 2026')])
    print('updated the title slide')

# ---- 3. insert the new slides (after-index refers to the original deck, and
#         nothing has been deleted yet, so the numbering still lines up) ------
for after_num, seq, layout, title, body in sorted(C.INSERT, key=lambda x: (x[0], x[1]), reverse=True):
    P.add_slide_after(prs, after_num - 1, layout, title, body)
print(f'inserted {len(C.INSERT)} slides')

# ---- 4. delete the obsolete slides. Their positions have shifted by however
#         many insertions landed before them, so adjust. ---------------------
inserted_after = sorted(item[0] for item in C.INSERT)
deleted = 0
for num in sorted(C.DELETE, reverse=True):
    shift = sum(1 for a in inserted_after if a < num)
    P.delete_slide(prs, num - 1 + shift)
    deleted += 1
print(f'deleted {deleted} slides')

prs.save(DST)
print(f'wrote {DST}: {original_count} -> {len(Presentation(DST).slides)} slides')

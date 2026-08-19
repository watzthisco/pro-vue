"""Helpers for rewriting slides in an existing PowerPoint deck."""
import copy

from pptx.enum.shapes import PP_PLACEHOLDER


TITLE_PH = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
SKIP_PH = {PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.DATE,
           PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.PICTURE}


def title_shape(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.type in TITLE_PH:
            return sh
    return None


def body_shape(slide):
    """The main content placeholder: the biggest text placeholder that isn't the
    title, the slide number, or a footer."""
    best = None
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.is_placeholder:
            continue
        t = sh.placeholder_format.type
        if t in TITLE_PH or t in SKIP_PH:
            continue
        area = (sh.width or 0) * (sh.height or 0)
        if best is None or area > (best.width or 0) * (best.height or 0):
            best = sh
    return best


def set_text_frame(shape, lines):
    """Replace a text frame's paragraphs with `lines`, a list of (level, text)."""
    tf = shape.text_frame
    # Keep the first paragraph so its list formatting survives, and clone it
    # for the rest.
    first = tf.paragraphs[0]
    template = copy.deepcopy(first._p)
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    for r in list(first.runs):
        r._r.getparent().remove(r._r)

    from pptx.text.text import _Paragraph
    paras = [first]
    for _ in lines[1:]:
        new_p = copy.deepcopy(template)
        for child in list(new_p):
            if child.tag.endswith('}r'):
                new_p.remove(child)
        first._p.getparent().append(new_p)
        paras.append(_Paragraph(new_p, tf))

    for para, (level, text) in zip(paras, lines):
        para.level = level
        if text:
            para.add_run().text = text


def set_slide(slide, title=None, body=None):
    if title is not None:
        sh = title_shape(slide)
        if sh is None:
            raise ValueError('slide has no title placeholder')
        set_text_frame(sh, [(0, title)])
    if body is not None:
        sh = body_shape(slide)
        if sh is None:
            raise ValueError('slide has no body placeholder')
        set_text_frame(sh, body)


def move_slide(prs, from_index, to_index):
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    el = ids[from_index]
    sldIdLst.remove(el)
    sldIdLst.insert(to_index, el)


def delete_slide(prs, index):
    sldIdLst = prs.slides._sldIdLst
    slide_id = list(sldIdLst)[index]
    rId = slide_id.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    sldIdLst.remove(slide_id)


def layout_by_name(prs, name):
    for layout in prs.slide_masters[0].slide_layouts:
        if layout.name == name:
            return layout
    raise KeyError(f'no layout named {name!r}')


def add_slide_after(prs, index, layout_name, title, body=None):
    """Append a new slide then move it into position after `index`."""
    layout = layout_by_name(prs, layout_name)
    slide = prs.slides.add_slide(layout)

    # Drop placeholders the caller isn't filling, so they don't render as
    # "Click to add text" prompts.
    for sh in list(slide.shapes):
        if not sh.is_placeholder:
            continue
        t = sh.placeholder_format.type
        if t in TITLE_PH:
            continue
        if t in SKIP_PH:
            continue
        if body is None:
            sh._element.getparent().remove(sh._element)

    set_slide(slide, title=title, body=body)
    move_slide(prs, len(prs.slides._sldIdLst) - 1, index + 1)
    return slide


def replace_graphic_with_text(slide, lines):
    """Swap a picture/diagram content placeholder for a text box in its place.

    Slides whose content is a screenshot of old code have nothing to rewrite,
    so the graphic is dropped and replaced with a text box at the same
    geometry.
    """
    from pptx.util import Pt

    target = None
    for sh in slide.shapes:
        if sh.is_placeholder and not sh.has_text_frame:
            t = sh.placeholder_format.type
            if t not in TITLE_PH and t not in SKIP_PH:
                target = sh
                break
    if target is None:
        return False

    left, top, width, height = target.left, target.top, target.width, target.height
    target._element.getparent().remove(target._element)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (level, text) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = level
        run = para.add_run()
        run.text = text
        run.font.size = Pt(18)
    return True
